"""
Text extraction from uploaded documents (.html, .docx, .pdf, .pptx).
Stateless — reads bytes and returns text; never stores files.
"""

MAX_CHARS = 5000


def extract_text(filename: str, content_bytes: bytes) -> str:
    """Dispatch extraction by file extension. Returns plain text, truncated to MAX_CHARS."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    extractors = {
        "html": _extract_html,
        "htm": _extract_html,
        "docx": _extract_docx,
        "pdf": _extract_pdf,
        "pptx": _extract_pptx,
    }
    fn = extractors.get(ext)
    if fn is None:
        raise ValueError(f"Unsupported file type: .{ext}")
    text = fn(content_bytes)
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS] + f"\n\n[Truncated — {len(text)} total characters]"
    return text


def _extract_html(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


def _extract_docx(content: bytes) -> str:
    import io
    from docx import Document

    doc = Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_pdf(content: bytes) -> str:
    import pymupdf

    doc = pymupdf.open(stream=content, filetype="pdf")
    pages = []
    for page in doc:
        text = page.get_text().strip()
        if text:
            pages.append(text)
    doc.close()
    return "\n\n--- Page Break ---\n\n".join(pages)


def _extract_pptx(content: bytes) -> str:
    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)
